'''
@modified: embedding_m_0714
2025.07.28 상태코드별 웹소켓 메세지 전송
2025.08.04 절대경로로 변경, page_content에 none값 방지
2025.08.14 chromadb 서버방식으로 변경, 배치 적재방식 삭제
2025.08.25 배치 적재방식 추가
'''
import asyncio # 표준 라이브러리 임포트
import os
import subprocess
import asyncpg # 써드파티 라이브러리 임포트
import nltk
import redis
from dotenv import load_dotenv
from fastapi import HTTPException
from pptx import Presentation
from langchain.chains.summarize import load_summarize_chain # LangChain 관련 임포트
from langchain_community.document_loaders import PyMuPDFLoader
from langchain_openai import OpenAIEmbeddings
from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
#from langchain_community.vectorstores import Chroma
from langchain_chroma import Chroma
import chromadb
from langchain_community.document_loaders import TextLoader
from langchain_openai import ChatOpenAI
from config import DATABASE_CONFIG, REDIS_CONFIG # 로컬 모듈 임포트
from datetime import datetime
import json  # JSON 처리를 위해 추가
from langchain.prompts import PromptTemplate
load_dotenv()
nltk.download('punkt')

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
CHROMA_HOST = os.getenv("CHROMA_HOST", "127.0.0.1")
CHROMA_PORT = int(os.getenv("CHROMA_PORT", "8002"))
CHROMA_CLIENT = chromadb.HttpClient(host=CHROMA_HOST, port=CHROMA_PORT)

redis_client = redis.Redis(**REDIS_CONFIG)

# 디렉토리 권한 확인 및 설정 함수
def ensure_directory_permissions(path, mode=0o775):
    """
    디렉토리 및 상위 디렉토리의 권한을 확인하고 설정합니다.
    """
    try:
        # 디렉토리가 이미 존재하는 경우
        if os.path.exists(path):
            current_mode = os.stat(path).st_mode & 0o777
            if current_mode != mode:
                print(f"권한 변경 필요: {path} (현재: {current_mode:o}, 목표: {mode:o})")
                os.chmod(path, mode)
                print(f"권한 변경 완료: {path}")
            else:
                print(f"권한 설정 정상: {path}")
            return
            
        # 디렉토리가 존재하지 않으면 상위 디렉토리 권한 확인
        parent = os.path.dirname(path)
        if parent and parent != path:  # 무한 루프 방지
            ensure_directory_permissions(parent, mode)
            
        # 디렉토리 생성 후 권한 설정
        os.makedirs(path, exist_ok=True)
        os.chmod(path, mode)
        print(f"디렉토리 생성 및 권한 설정 완료: {path}")
    except Exception as e:
        print(f"권한 설정 중 오류 발생: {path} - {e}")
        import traceback
        print(traceback.format_exc())

async def get_file_name_from_db(file_id: str) -> str:
    query = """
        SELECT file_nm
        FROM cls_file_mst
        WHERE file_id = $1
    """
    
    try:
        print(f"Fetching file name for file_id: {file_id}")  # 디버깅 로그
        conn = await asyncpg.connect(**DATABASE_CONFIG)
        try:
            result = await conn.fetchval(query, file_id)
            print(f"Query result for file_id {file_id}: {result}")  # 디버깅 로그
            return result
        finally:
            await conn.close()
    except Exception as e:
        print(f"Error fetching file name from DB: {e}")  # 디버깅 로그
        raise HTTPException(status_code=500, detail=f"Database error: {str(e)}")

def load_text_file(file_path, title=None, file_type_str=None):
    """
    텍스트 파일 로드 및 처리
    """
    loader = TextLoader(file_path)
    documents = loader.load()

    # 메타데이터만 추가
    for doc in documents:
        # page_content가 None이거나 빈 문자열인 경우 건너뛰기
        if not doc.page_content or doc.page_content.strip() == "":
            continue
            
        doc.metadata = {
            "title": title if title else "",
            "종류": file_type_str if file_type_str else ""
        }
    
    # 유효한 문서만 반환
    return [doc for doc in documents if doc.page_content and doc.page_content.strip()]

def load_pdf_file(file_path, title=None, file_type_str=None):
    """
    PDF 파일 로드 및 처리 - 모든 페이지를 하나의 문서로 합침
    """
    loader = PyMuPDFLoader(file_path)
    documents = loader.load()
    
    # 유효한 페이지 내용만 필터링
    valid_contents = []
    for doc in documents:
        if doc.page_content and doc.page_content.strip():
            valid_contents.append(doc.page_content)
    
    # 유효한 내용이 없는 경우 빈 문자열로 처리
    if not valid_contents:
        combined_text = ""
    else:
        combined_text = "\n\n".join(valid_contents)
    
    # 하나의 Document로 만듦
    metadata = {
        "title": title if title else "",
        "종류": file_type_str if file_type_str else "",
        "total_pages": len(documents)  # 전체 페이지 수 정보는 유지
    }
    
    return [Document(page_content=combined_text, metadata=metadata)]

def load_ppt_file(file_path, title=None, file_type_str=None):
    """
    파워포인트 파일 로드 및 텍스트 추출
    """
    def extract_text_from_pptx(file_path):
        presentation = Presentation(file_path)
        texts = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
        return texts
    
    # 파일 확장자가 .ppt일 경우 변환 후 사용
    if file_path.endswith(".ppt"):
        file_path = convert_ppt_to_pptx(file_path)

    texts = extract_text_from_pptx(file_path)
    documents = []
    for text in texts:
        # 빈 텍스트는 건너뛰기
        if not text or text.strip() == "":
            continue
            
        metadata = {
            "title": title if title else "",
            "종류": file_type_str if file_type_str else ""
        }
        document = Document(
            page_content=text,
            metadata=metadata
        )
        documents.append(document)

    return documents

def convert_ppt_to_pptx(ppt_path):
    """LibreOffice를 사용하여 PPT를 PPTX로 변환"""
    output_path = ppt_path.replace(".ppt", ".pptx")
    
    command = [
        "libreoffice", "--headless", "--convert-to", "pptx", ppt_path, "--outdir", os.path.dirname(ppt_path)
    ]
    subprocess.run(command, stdout=subprocess.PIPE, stderr=subprocess.PIPE)

    if os.path.exists(output_path):
        return output_path  # 변환 성공
    else:
        raise FileNotFoundError(f"변환 실패: {ppt_path}")
    
async def load_file(file_path, file_type_cd, file_id):
    """
    파일을 로드하고 태그를 추가하는 함수
    """
    type_mapping = {
        "S": "강의일정",
        "N": "공지사항",
        "M": "수업자료",
        "V": "강의영상"
    }
    file_type_str = type_mapping.get(file_type_cd, "기타")
    
    # DB에서 파일명 가져오기
    title = await get_file_name_from_db(file_id)
    if not title:
        print(f"Database did not return a title for file_id {file_id}. Using file name from path.")
        title = os.path.basename(file_path)
    
    print(f"Loading file with title: {title} and type: {file_type_str}")  # 디버깅 로그 추가

    ext = os.path.splitext(file_path)[-1].lower()
    
    if ext == ".txt":
        documents = load_text_file(file_path, title, file_type_str)
    elif ext in [".pptx", ".ppt"]: 
        documents = load_ppt_file(file_path, title, file_type_str)
    elif ext == ".pdf":
        documents = load_pdf_file(file_path, title, file_type_str)
    else:
        raise ValueError(f"FILE_FORMAT_ERROR: 지원되지 않는 파일 형식입니다: {ext}")

    return documents

async def update_status_in_db(status, dt_tag, file_id, cls_id=None, user_id=None):
    """
    데이터베이스에서 파일 상태 업데이트 및 웹소켓 알림 전송
    """
    if dt_tag == 'NO_UPDATE':
        query = """
            UPDATE cls_file_mst
            SET upload_status = $1
            WHERE FILE_ID = $2
        """
    else:
        query = f"""
            UPDATE cls_file_mst
            SET upload_status = $1,
            {dt_tag} = CURRENT_TIMESTAMP AT TIME ZONE 'Asia/Seoul'
            WHERE FILE_ID = $2
        """

    try:
        conn = await asyncpg.connect(**DATABASE_CONFIG)
        try:
            await conn.execute(query, status, file_id)
            print(f"DB 업데이트: {file_id} - {status}")
            
            # 웹소켓 알림 전송
            if cls_id and user_id:
                stream_key = f"FDQ:{cls_id}:{user_id}"        
                redis_hash_key = f"{cls_id}:material_status"
                
                # 상태에 따라 적절한 메시지 생성
                status_message = ""
                task_type = "status_update"
                
                if status == "EP02":
                    status_message = "튜터학습중"
                    task_type = "embedding"
                    update_hash_and_notify_task(redis_hash_key, stream_key, task_type, file_id, {"embedding_status": "processing", "message": status_message})
                elif status == "EP03":
                    status_message = "튜터학습완료"
                    task_type = "embedding"
                    update_hash_and_notify_task(redis_hash_key, stream_key, task_type, file_id, {"embedding_status": "completed", "message": status_message})
                elif status == "EP04":
                    status_message = "튜터학습실패"
                    task_type = "embedding"
                    update_hash_and_notify_task(redis_hash_key, stream_key, task_type, file_id, {"embedding_status": "failed", "message": status_message})
                elif status == "SM02":
                    status_message = "주차별학습중"
                    task_type = "summary"
                    update_hash_and_notify_task(redis_hash_key, stream_key, task_type, file_id, {"summary_status": "processing", "message": status_message})
                elif status == "SM03":
                    status_message = "주차별학습완료"
                    task_type = "summary"
                    update_hash_and_notify_task(redis_hash_key, stream_key, task_type, file_id, {"summary_status": "completed", "message": status_message})
                elif status == "SM04":
                    status_message = "주차별학습실패"
                    task_type = "summary"
                    update_hash_and_notify_task(redis_hash_key, stream_key, task_type, file_id, {"summary_status": "failed", "message": status_message})
                elif status == "FU04":
                    status_message = "업로드 실패"
                    task_type = "upload"
                    update_hash_and_notify_task(redis_hash_key, stream_key, task_type, file_id, {"upload_status": "failed", "message": status_message})
                elif status == "EP05":
                    status_message = "임베딩 처리 불가"
                    task_type = "embedding"
                    update_hash_and_notify_task(redis_hash_key, stream_key, task_type, file_id, {"embedding_status": "unsupported", "message": status_message})
                elif status == "FD01":
                    status_message = "삭제중"
                    task_type = "delete"
                    update_hash_and_notify_task(redis_hash_key, stream_key, task_type, file_id, {"delete_status": "processing", "message": status_message})
                    
        finally:
            await conn.close()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Database error: {str(e)}")

# 주차 정보 추출 (week_num 추출)
def extract_week_info(file_path: str) -> int:
    """
    파일 경로에서 주차 정보 추출
    예: /uploaded_files/CLASS_ID/week_1/file.pdf -> 1
    """
    try:
        path_parts = file_path.split('/')
        for part in path_parts:
            if part.lower().startswith('week'):
                week_num = ''.join(filter(str.isdigit, part))
                return int(week_num) if week_num else 0
        return 0
    except Exception:
        return 0

def preprocess_text_with_metadata(text: str, metadata: dict) -> str:
    """
    텍스트에 메타데이터 정보를 추가
    """
    week = metadata.get('week', '알 수 없음')
    file_type = metadata.get('종류', '알 수 없음')
    title = metadata.get('title', '알 수 없음')
    
    prefix = f"이 문서는 {week}주차에 교수자가 [{file_type}] 카테고리에 업로드한 내용입니다. "
    prefix += f"참고한 파일의 제목은 [{title}]입니다.\n\n"
    
    return prefix + text

async def identify_new_files(cls_id: str) -> list:
    """
    임베딩되지 않은 새로운 파일들을 식별
    """
    query = """
        SELECT 
            file_id,
            cls_id,
            file_type_cd,
            file_nm,
            file_path,
            upload_status
        FROM cls_file_mst
        WHERE cls_id = $1
        AND upload_status IN ('FU03')  -- 파일 업로드 완료인 파일 찾기
        AND file_ext <> 'mp4' -- 영상 파일 제외
    """
        #     AND file_del_req_dt IS NULL
        # AND emb_del_comp_dt IS NULL
    #         AND upload_status IN ('FU03', 'EP04')  -- 파일 업로드 완료이거나 임베딩 처리 실패

    try:
        conn = await asyncpg.connect(**DATABASE_CONFIG)
        try:
            files = await conn.fetch(query, cls_id)
            return [dict(file) for file in files]
        finally:
            await conn.close()
    except Exception as e:
        print(f"Error identifying new files: {e}")
        return []

async def update_summaries_for_newly_embedded_files(cls_id: str, file_path: str, docs: list, user_id: str = None):
    try:
        # 주차 정보를 file_path에서 추출하는 대신 docs의 메타데이터에서 가져옵니다
        # 모든 docs는 같은 파일에서 왔으므로 첫 번째 doc의 week 값을 사용합니다
        if not docs or len(docs) == 0:
            print("No documents to process")
            return
            
        # page_content가 None이거나 빈 문자열인 문서 필터링
        docs = [doc for doc in docs if doc.page_content and doc.page_content.strip()]
        
        if not docs or len(docs) == 0:
            print("No valid documents after filtering")
            return
            
        # 파일 ID 추출
        file_id = docs[0].metadata.get("file_id", "")
        if not file_id:
            print("No file_id found in metadata")
            return
            
        # 요약 시작 상태 업데이트
        await update_status_in_db("SM02", 'SUMRY_START_DT', file_id, cls_id, user_id)
        print(f"Summary process started for file_id: {file_id}")
            
        week = docs[0].metadata.get("week", 0)
        if week <= 0:
            # 메타데이터에서 유효한 주차 정보를 찾지 못한 경우 DB에서 조회
            try:
                if file_id:
                    conn = await asyncpg.connect(**DATABASE_CONFIG)
                    try:
                        query = """
                            SELECT week_num
                            FROM cls_file_mst
                            WHERE file_id = $1
                        """
                        week = await conn.fetchval(query, file_id)
                        week = int(week) if week is not None else 0
                        print(f"Retrieved week {week} from database for file_id: {file_id}")
                    finally:
                        await conn.close()
            except Exception as db_err:
                print(f"Error fetching week from database: {db_err}")
                # 요약 실패 상태 업데이트
                await update_status_in_db("SM04", 'SUMRY_FAIL_DT', file_id, cls_id, user_id)
                print(f"Summary process failed for file_id: {file_id}")
                return
                
        if week <= 0:
            print("No valid week information found in metadata or database")
            # 요약 실패 상태 업데이트
            await update_status_in_db("SM04", 'SUMRY_FAIL_DT', file_id, cls_id, user_id)
            print(f"Summary process failed for file_id: {file_id} - No valid week information")
            return
            
        print(f"Creating summary for week: {week}")
            
        # ChromaDB에서 같은 주차의 모든 문서 검색
        #chroma_db_path = os.path.join("/app/tutor/db/chromadb", cls_id)
        embeddings = OpenAIEmbeddings(model="text-embedding-3-large")
        db = Chroma(
            #persist_directory=chroma_db_path, 
            client=CHROMA_CLIENT, # 서버방식으로 변경
            embedding_function=embeddings, 
            collection_name=cls_id
            )
        
        # 현재 주차의 모든 문서 검색
        results = db.get(
            where={"week": week}
        )
        
        if not results or not results['documents']:
            print(f"No documents found in ChromaDB for week {week}")
            documents = docs  # 현재 처리 중인 문서만 사용
        else:
            # ChromaDB의 문서와 현재 문서 결합
            documents = []
            for i, content in enumerate(results['documents']):
                # page_content가 None이거나 빈 문자열인 경우 건너뛰기
                if not content or content.strip() == "":
                    continue
                    
                metadata = results['metadatas'][i]
                doc = Document(
                    page_content=content,
                    metadata=metadata
                )
                documents.append(doc)
            documents.extend(docs)  # 현재 처리 중인 문서 추가
            
        print(f"Total {len(documents)} documents found for week {week}")
        
        # 토큰 수 추정 함수
        def estimate_tokens(text):
            """
            텍스트의 대략적인 토큰 수를 추정합니다.
            """
            # 영어는 약 4자당 1토큰, 한국어는 약 2자당 1토큰으로 추정
            korean_chars = sum(1 for char in text if '\u3131' <= char <= '\u318E' or '\uAC00' <= char <= '\uD7A3')
            other_chars = len(text) - korean_chars
            estimated_tokens = (korean_chars / 2) + (other_chars / 4)
            return int(estimated_tokens)
        
        # 토큰 제한을 위한 문서 크기 제한 함수
        def limit_documents_for_tokens(docs, max_tokens_per_doc=8000, max_total_docs=100):
            """
            토큰 제한을 고려하여 문서 수와 크기를 제한합니다.
            """
            limited_docs = []
            total_tokens = 0
            max_tokens_total = 250000  # 안전한 토큰 제한 (300,000보다 여유있게)
            
            for doc in docs[:max_total_docs]:  # 최대 문서 수 제한
                content = doc.page_content
                estimated_tokens = estimate_tokens(content)
                
                if estimated_tokens > max_tokens_per_doc:
                    # 문서가 너무 크면 앞부분만 사용
                    # 토큰 수에 맞게 텍스트 자르기
                    target_chars = max_tokens_per_doc * 2  # 안전한 추정
                    content = content[:target_chars] + "..."
                    doc.page_content = content
                    estimated_tokens = estimate_tokens(content)
                
                limited_docs.append(doc)
                total_tokens += estimated_tokens
                
                # 전체 토큰 수 제한
                if total_tokens > max_tokens_total:
                    print(f"토큰 제한에 도달했습니다. {len(limited_docs)}개 문서, 총 {total_tokens} 토큰")
                    break
            
            print(f"토큰 제한 적용: {len(limited_docs)}개 문서, 총 {total_tokens} 토큰")
            return limited_docs
        
        # 문서 크기 제한 적용
        limited_documents = limit_documents_for_tokens(documents)
        
        # 요약을 위한 프롬프트 템플릿 정의 (간결하게 수정)
        map_prompt_template = """대학 강의 수업 자료입니다. 핵심 개념, 중요 용어, 주요 내용을 간결하게 요약해주세요.

{text}

요약:"""
        map_prompt = PromptTemplate(template=map_prompt_template, input_variables=["text"])

        combine_prompt_template = """대학 강의 주차별 수업 자료 요약입니다. 아래 요약들을 통합하여 핵심 내용을 3-4개 섹션으로 구성해 주세요.

{text}

종합 요약:"""
        combine_prompt = PromptTemplate(template=combine_prompt_template, input_variables=["text"])

        # 요약 생성
        try:
            llm = ChatOpenAI(temperature=0, max_tokens=2000)  # max_tokens 제한 추가
            chain = load_summarize_chain(
                llm,
                chain_type="map_reduce",
                map_prompt=map_prompt,
                combine_prompt=combine_prompt,
                verbose=False
            )
            chain_result = chain.invoke(limited_documents)
            summary = chain_result.get('output_text', '') if isinstance(chain_result, dict) else str(chain_result)
            
            # 타이틀 생성 (한국어)
            title_prompt = PromptTemplate(
                template="다음 대학 강의 요약의 핵심 주제를 반영한 간결한 제목을 한 줄로 작성해주세요:\n\n{text}\n\n제목:",
                input_variables=["text"]
            )
            title_result = llm.invoke(title_prompt.format(text=summary[:500]))
            title = title_result.content.strip() if hasattr(title_result, 'content') else str(title_result)
        except Exception as llm_err:
            print(f"Error generating summary with LLM: {llm_err}")
            # 요약 실패 상태 업데이트
            await update_status_in_db("SM04", 'SUMRY_FAIL_DT', file_id, cls_id, user_id)
            print(f"Summary process failed for file_id: {file_id}")
            return
        
        # JSON 형식으로 데이터 구성
        content_data = json.dumps({
            "title": title,
            "content": summary
        }, ensure_ascii=False)
        
        # 현재 시간 가져오기
        current_time = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        # DB에 요약 저장 (UPDATE 또는 INSERT)
        query = """
            WITH upsert AS (
                UPDATE cls_weekly_learn 
                SET content_data = $3,
                    upd_dt = CURRENT_TIMESTAMP AT TIME ZONE 'Asia/Seoul'
                WHERE cls_id = $1 
                AND week_num = $2
                RETURNING *
            )
            INSERT INTO cls_weekly_learn (cls_id, week_num, content_data, ins_dt, upd_dt)
            SELECT $1, $2, $3, CURRENT_TIMESTAMP AT TIME ZONE 'Asia/Seoul', CURRENT_TIMESTAMP AT TIME ZONE 'Asia/Seoul'
            WHERE NOT EXISTS (SELECT * FROM upsert);
        """
        
        try:
            conn = await asyncpg.connect(**DATABASE_CONFIG)
            try:
                await conn.execute(
                    query, 
                    cls_id,    # $1
                    week,      # $2
                    content_data  # $3
                )
                print(f"Week {week} summary updated successfully")
                
                # 요약 완료 상태 업데이트
                await update_status_in_db("SM03", 'SUMRY_COMP_DT', file_id, cls_id, user_id)
                print(f"Summary process completed for file_id: {file_id}")
            finally:
                await conn.close()
        except Exception as db_err:
            print(f"Error updating summary in database: {db_err}")
            # 요약 실패 상태 업데이트
            await update_status_in_db("SM04", 'SUMRY_FAIL_DT', file_id, cls_id, user_id)
            print(f"Summary process failed for file_id: {file_id}")
            return
        
        return f"Updated summary for week {week}"
        
    except Exception as e:
        print(f"Error updating summaries: {e}")
        import traceback
        print(traceback.format_exc())
        
        # 파일 ID가 있는 경우 실패 상태 업데이트
        if docs and len(docs) > 0:
            file_id = docs[0].metadata.get("file_id", "")
            if file_id:
                await update_status_in_db("SM04", 'SUMRY_FAIL_DT', file_id, cls_id, user_id)
                print(f"Summary process failed for file_id: {file_id}")
        
        raise

# async def process_and_embed_files(cls_id: str, files_info: list, notify_filestatus_toclients):
async def process_and_embed_files(cls_id: str, files_info: list, user_id: str):
    """
    파일 처리 및 임베딩 수행
    각 파일에 대해 임베딩 작업 진행 상황을 Redis 해시와 스트림에 기록
    """
    base_path = f"/app/tutor/download_files/{cls_id}"

    # 디렉토리 권한 확인 및 설정
    ensure_directory_permissions(base_path)
  
    # 기본 디렉토리 권한 설정
    os.chmod(base_path, 0o775)
          
    # Redis 해시와 스트림 키 설정
    redis_hash_key = f"{cls_id}:material_status"
    stream_key = f"FDQ:{cls_id}:{user_id}"

    for file_info in files_info:
        file_id = file_info["file_id"]
        file_path = file_info["file_path"]  # DB에서 가져온 전체 경로 사용
        file_type_cd = file_info["file_type_cd"]

        update_hash_and_notify_task(redis_hash_key, stream_key, "embedding", file_id, {"embedding_status": "start"})

        try:
            # 파일 존재 및 확장자 확인
            if not os.path.exists(file_path):
                print(f"File not found: {file_path}")
                update_hash_and_notify_task(redis_hash_key, stream_key, "embedding", file_id, {"embedding_status": "failed", "error": "파일을 찾을 수 없습니다."})                
                # 파일 업로드 실패 코드로 설정, dwld_fail_dt에 시간값 넣기
                await update_status_in_db("FU04", 'DWLD_FAIL_DT', file_id, cls_id, user_id)
                redis_client.hdel(redis_hash_key, file_id)
                continue

            _, ext = os.path.splitext(file_path)
            if not ext:
                print(f"No file extension found for: {file_path}")
                update_hash_and_notify_task(redis_hash_key, stream_key, "embedding", file_id, {"embedding_status": "failed", "error": "파일 확장자가 없습니다."})                
                # 임베딩 처리 불가 코드로 설정, 시간값 업데이트 없음
                await update_status_in_db("EP05", 'NO_UPDATE', file_id, cls_id, user_id)
                redis_client.hdel(redis_hash_key, file_id)
                continue

            # 파일 로드 및 태그 추가
            try:
                documents = await load_file(file_path, file_type_cd, file_id)
                await update_status_in_db("EP02", 'EMB_START_DT', file_id, cls_id, user_id)
                print(f"{file_id} 파일 로드 성공!")
            except ValueError as ve:
                if "지원되지 않는 파일 형식" in str(ve):
                    # 파일 형식 오류는 임베딩 처리 불가로 기록
                    print(f"File format error for {file_id}: {str(ve)}")
                    update_hash_and_notify_task(redis_hash_key, stream_key, "embedding", file_id, 
                                    {"embedding_status": "failed", "error": str(ve)})
                    await update_status_in_db("EP05", 'NO_UPDATE', file_id, cls_id, user_id)
                    redis_client.hdel(redis_hash_key, file_id)
                    continue
                else:
                    # 다른 ValueError는 다시 던져서 임베딩 실패로 처리
                    raise

            # 텍스트 분할 및 메타데이터 업데이트 (토큰 제한 고려)
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=2000,  # 청크 크기 2000자로 설정
                chunk_overlap=200,
                length_function=len,
                separators=["\n\n", "\n", ".", " ", ""]  # 자연스러운 문단 구분을 위한 구분자
            )
            docs = text_splitter.split_documents(documents)
            
            # page_content가 None이거나 빈 문자열인 문서 필터링
            docs = [doc for doc in docs if doc.page_content and doc.page_content.strip()]
            
            # 임베딩용 문서 수 제한 없음 (배치 처리 사용)
            print(f"임베딩용 문서 수: {len(docs)}개")
            
            today = datetime.now().strftime("%Y-%m-%d")
            for doc in docs:
                # page_content가 None이거나 빈 문자열인 경우 건너뛰기
                if not doc.page_content or doc.page_content.strip() == "":
                    continue
                    
                # 메타데이터 간소화
                doc.metadata = {
                    "week": extract_week_info(file_path),
                    "종류": documents[0].metadata.get("종류", "Unknown"),
                    "title": documents[0].metadata.get("title", "Unknown"),
                    "embedding_date": today,
                    "file_id": file_id
                }
                doc.page_content = preprocess_text_with_metadata(doc.page_content, doc.metadata)

            # 임베딩 생성 및 저장 (배치 처리)
            try:
                embeddings = OpenAIEmbeddings(model="text-embedding-3-large")
                
                # 배치 크기 설정
                batch_size = 50
                total_docs = len(docs)
                
                print(f"배치 처리 시작: 총 {total_docs}개 문서를 {batch_size}개씩 처리")
                
                # Retry logic for dimension mismatch
                max_retries = 1
                for attempt in range(max_retries + 1):
                    try:
                        vector_store = Chroma(
                            client=CHROMA_CLIENT,
                            embedding_function=embeddings,
                            collection_name=cls_id
                        )
                        
                        # 배치별로 문서 처리
                        for i in range(0, total_docs, batch_size):
                            batch_docs = docs[i:i + batch_size]
                            batch_num = (i // batch_size) + 1
                            total_batches = (total_docs + batch_size - 1) // batch_size
                            
                            print(f"배치 {batch_num}/{total_batches} 처리 중... ({len(batch_docs)}개 문서)")
                            
                            try:
                                vector_store.add_documents(batch_docs)
                                print(f"배치 {batch_num} 완료")
                            except Exception as batch_e:
                                # Dimension mismatch inside batch
                                if "dimension" in str(batch_e) and "1536" in str(batch_e) and "3072" in str(batch_e):
                                    raise batch_e
                                
                                print(f"배치 {batch_num} 처리 중 오류: {batch_e}")
                                # 개별 문서 재시도
                                for j, doc in enumerate(batch_docs):
                                    try:
                                        vector_store.add_documents([doc])
                                        print(f"개별 문서 {i+j+1} 처리 완료")
                                    except Exception as doc_e:
                                        if "dimension" in str(doc_e) and "1536" in str(doc_e) and "3072" in str(doc_e):
                                            raise doc_e
                                        print(f"개별 문서 {i+j+1} 처리 실패: {doc_e}")
                                        continue
                        
                        print(f"모든 문서 임베딩 완료: {total_docs}개")
                        break # Success, exit retry loop

                    except Exception as e:
                        # Catch dimension mismatch and retry
                        if "dimension" in str(e) and "1536" in str(e) and "3072" in str(e):
                            if attempt < max_retries:
                                print(f"Warning: Dimension mismatch detected for {cls_id}. Deleting and recreating collection... (Attempt {attempt+1})")
                                CHROMA_CLIENT.delete_collection(cls_id)
                                continue
                            else:
                                print(f"Error: Dimension mismatch persists after retry for {cls_id}.")
                                raise e
                        else:
                            raise e

            except Exception as chroma_e:
                print(f"ChromaDB 오류: {chroma_e}")
                raise

            # 임베딩 성공 처리 및 요약문 작성
            if docs and len(docs) > 0 and docs[0].metadata["week"] > 0:
                # 먼저 임베딩 완료 상태 업데이트
                await update_status_in_db("EP03", 'EMB_COMP_DT', file_id, cls_id, user_id)
                
                # 요약문 작성 (내부에서 SM03 상태까지 업데이트됨)
                print(f"{file_id} 파일 임베딩 완료, 요약 작성 시작...")
                # 조교 서비스에서는 요약 기능 미사용으로 주석 처리
                # await update_summaries_for_newly_embedded_files(cls_id, file_path, docs, user_id)
                
                update_hash_and_notify_task(redis_hash_key, stream_key, "embedding", file_id, {"embedding_status": "completed"})
                print(f"{file_id} 파일 임베딩 및 요약 생성 성공!")
            else:
                # 주차 정보가 없는 경우에는 요약 작성 없이 DB 업데이트
                # 주차 정보가 없는 경우에는 요약 작성 없이 DB 업데이트
                await update_status_in_db("EP03", 'EMB_COMP_DT', file_id, cls_id, user_id)
                update_hash_and_notify_task(redis_hash_key, stream_key, "embedding", file_id, {"embedding_status": "completed"})
                print(f"{file_id} 파일 임베딩 성공! (주차 정보 없음)")
            
            redis_client.hdel(redis_hash_key, file_id)

        except Exception as e:
            print(f"Error processing file {file_id}: {str(e)}")
            import traceback
            print(traceback.format_exc())
            
            # 예외 메시지에 따라 파일 처리 실패 또는 임베딩 실패로 구분
            if "지원되지 않는 파일 형식" in str(e) or "No file extension" in str(e):
                # 파일 형식 관련 오류는 임베딩 처리 불가로 기록 (시간값 업데이트 없음)
                await update_status_in_db("EP05", 'NO_UPDATE', file_id, cls_id, user_id)
                error_message = "파일 형식 오류 (임베딩 처리 불가): " + str(e)
            else:
                # 그 외 오류는 임베딩 실패로 기록
                await update_status_in_db("EP04", 'EMB_FAIL_DT', file_id, cls_id, user_id)
                error_message = "임베딩 실패: " + str(e)
                
            update_hash_and_notify_task(redis_hash_key, stream_key, "embedding", file_id,
                              {"embedding_status": "failed", "error": error_message},
                              {"status": "failed", "error": error_message})
            
            redis_client.hdel(redis_hash_key, file_id)
            print(f"{file_id} 처리 실패: {error_message}")





# async def process_new_files_and_update_summaries(cls_id: str, notify_filestatus_toclients):
async def process_new_files_and_update_summaries(cls_id: str, user_id):
    """
    새 파일 처리 및 주차별 요약 업데이트를 수행하는 메인 함수
    """
    try:
        # 1. 새로운 파일 식별 (FU03, EP04 상태의 파일)
        new_files = await identify_new_files(cls_id)
        if not new_files:
            print("No new files to process")
            return
            
        print(f"Found {len(new_files)} files to process")
        
        # 2. 파일 정보 처리 및 임베딩 진행
        files_info = []
        for file in new_files:
            file_info = {
                "file_name": file['file_nm'],
                "file_id": file['file_id'],
                "file_type_cd": file['file_type_cd'],
                "file_path": file['file_path']
            }
            files_info.append(file_info)
            print(f"Processing file: {file_info}")
        
        # 3. 임베딩 진행
        # await process_and_embed_files(cls_id, files_info, notify_filestatus_toclients)
        await process_and_embed_files(cls_id, files_info, user_id)
        print("Processing completed successfully")
        
    except Exception as e:
        print(f"Error in main processing: {str(e)}")
        raise e


def update_hash_and_notify_task(redis_hash_key: str, stream_key: str, task_type:str, file_id: str, update_fields: dict, extra_fields: dict = None):
    try:
        existing = redis_client.hget(redis_hash_key, file_id)
        status_info = {}
        if existing:
            status_info = json.loads(existing.decode('utf-8'))
        status_info.update(update_fields)
        redis_client.hset(redis_hash_key, file_id, json.dumps(status_info))
        message = {
            "task_type": task_type,
            "file_id": file_id,
            "hash_info": json.dumps(status_info)
        }
        if extra_fields:
            message.update(extra_fields)
        redis_client.xadd(stream_key, message, id='*')
    except Exception as e:
        print(f"Redis update error for {file_id}: {e}")

# 예제 실행
if __name__ == "__main__":
    json_data = {
        "cls_id": "2024-20-903102-2-01"
        ,"user_id": "45932"
    }

    # notify_filestatus_toclients를 동기 함수로 처리
    # asyncio.run(process_new_files_and_update_summaries(json_data["cls_id"], lambda x: print(x)))
    asyncio.run(process_new_files_and_update_summaries(json_data["cls_id"],json_data["user_id"]))




# 정리 : raw_file 이 경로에 들어와 있고, FU03(업로드 완료 된 상태)일때 임베딩 -> 새롭게 임베딩이 된 파일이 속한 주차에 대해서는 요약문 다시 만듦
# 기존 파일 임베딩 처리 후 요약문 다시 만들기 추가